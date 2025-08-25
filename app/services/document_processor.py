import os
import uuid
from typing import List, Dict, Any
from pathlib import Path
import logging

from langchain.text_splitter import (
    RecursiveCharacterTextSplitter,
    CharacterTextSplitter,
    SentenceTransformersTokenTextSplitter
)
from langchain.docstore.document import Document as LangchainDocument
import fitz
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation

from app.config import ChunkingStrategy, ChunkingConfig
from app.model_schemas import Document, DocumentStatus

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf,
            '.txt': self._extract_txt,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx
        }

    def extract_text(self, file_path: str) -> str:
        """Extract text from various file formats."""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        try:
            return self.supported_formats[file_extension](file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise

    def _extract_pdf(self, file_path: str) -> str:
      """Extract text from PDF files, preserving structure and links."""
      structured_content = []
      pdf_document = fitz.open(file_path)
      for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        blocks = page.get_text("dict")["blocks"]  # Extract text blocks with structure
        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    line_text = ""
                    link_uri = None  # Track the link for this line, if any
                    for span in line["spans"]:
                        span_text = span["text"]
                        line_text += span_text
                        # Check if the span overlaps with any link
                        for link in page.get_links():
                            if "uri" in link:
                                link_rect = fitz.Rect(link["from"])
                                span_rect = fitz.Rect(span["bbox"])
                                if span_rect.intersects(link_rect):
                                    link_uri = link["uri"]  # Store the link URI
                    # Add the line text to the structured content
                    structured_content.append(line_text)
                    # Add the link (if any) after the line text
                    if link_uri:
                        structured_content.append(link_uri)
      return "\n".join(structured_content).strip()

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT files."""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX files."""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    def chunk_text(self, text: str, config: ChunkingConfig) -> List[LangchainDocument]:
        """Chunk text based on the specified strategy."""
        if config.strategy == ChunkingStrategy.FIXED_SIZE:
            splitter = CharacterTextSplitter(
                chunk_size=config.chunk_size,
                chunk_overlap=config.chunk_overlap,
                separator="\n\n"
            )
        elif config.strategy == ChunkingStrategy.RECURSIVE_TEXT:
            separators = config.separators or ["\n\n", "\n", " ", ""]
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=config.chunk_size,
                chunk_overlap=config.chunk_overlap,
                separators=separators
            )
        elif config.strategy == ChunkingStrategy.SENTENCE:
            splitter = SentenceTransformersTokenTextSplitter(
                chunk_overlap=config.chunk_overlap,
                tokens_per_chunk=config.chunk_size
            )
        elif config.strategy == ChunkingStrategy.SEMANTIC:
            # For now, use recursive as semantic chunking requires more complex implementation
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=config.chunk_size,
                chunk_overlap=config.chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
        else:
            raise ValueError(f"Unsupported chunking strategy: {config.strategy}")

        return splitter.create_documents([text])

    def process_document(
        self, 
        file_path: str, 
        filename: str, 
        configuration_name: str,
        chunking_config: ChunkingConfig,
        metadata: Dict[str, Any] = None
    ) -> Document:
        """Process a document: extract text, create document object."""
        try:
            # Extract text
            content = self.extract_text(file_path)
            
            # Get file info
            file_size = os.path.getsize(file_path)
            file_type = Path(file_path).suffix.lower().lstrip('.')
            
            # Create document
            document = Document(
                id=str(uuid.uuid4()),
                filename=filename,
                content=content,
                metadata=metadata or {},
                status=DocumentStatus.UPLOADED,
                configuration_name=configuration_name,
                file_size=file_size,
                file_type=file_type
            )
            
            logger.info(f"Successfully processed document: {filename}")
            return document
            
        except Exception as e:
            logger.error(f"Error processing document {filename}: {str(e)}")
            raise

    def get_chunks(self, document: Document, chunking_config: ChunkingConfig) -> List[LangchainDocument]:
        """Get chunks from a processed document."""
        chunks = self.chunk_text(document.content, chunking_config)
        
        # Add metadata to chunks
        for i, chunk in enumerate(chunks):
            chunk.metadata.update({
                "document_id": document.id,
                "filename": document.filename,
                "chunk_index": i,
                "configuration_name": document.configuration_name,
                **document.metadata
            })
        
        return chunks
        
    def process_text_content(
        self, 
        content: str, 
        filename: str, 
        configuration_name: str,
        chunking_config: ChunkingConfig,
        metadata: Dict[str, Any] = None
    ) -> Document:
        """Process raw text content without requiring a file.
        
        Args:
            content: The raw text content to process
            filename: A name to identify the document
            configuration_name: The configuration to use
            chunking_config: Configuration for text chunking
            metadata: Optional metadata for the document
            
        Returns:
            Document object with the processed content
        """
        try:
            # Create document with the provided text content
            document = Document(
                id=str(uuid.uuid4()),
                filename=filename,
                content=content,
                metadata=metadata or {},
                status=DocumentStatus.UPLOADED,
                configuration_name=configuration_name,
                file_size=len(content.encode('utf-8')),  # Size in bytes
                file_type='txt'  # Default to txt type for raw text
            )
            
            logger.info(f"Successfully processed text content as document: {filename}")
            return document
            
        except Exception as e:
            logger.error(f"Error processing text content as document {filename}: {str(e)}")
            raise
